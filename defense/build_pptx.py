"""Build pptx presentation from markdown slide files.

Reads defense/slides/*.md in lexicographic order. Each MD file becomes one slide.
No styling applied — title, body text, images, and notes are inserted as-is.
Tables are flattened to bulleted/numbered lists; markdown bullets become '○'.
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

DEFENSE_DIR = Path(__file__).resolve().parent
SLIDES_DIR = DEFENSE_DIR / "slides"
OUTPUT_PATH = DEFENSE_DIR / "presentation.pptx"

SLIDE_W_EMU = 12192000  # 13.333"  → 16:9 widescreen
SLIDE_H_EMU = 6858000   # 7.5"

ID_PATTERN = re.compile(r"^(?:\d+|[A-Z]+-\d+[a-z]?|SC-[A-Z])$")


def split_sections(md_text: str) -> dict[int, str]:
    sections: dict[int, str] = {}
    for chunk in re.split(r"\n---\n", md_text):
        chunk = chunk.strip()
        if not chunk:
            continue
        m = re.match(r"##\s*(\d+)\.\s*[^\n]*\n+(.*)", chunk, re.DOTALL)
        if m:
            sections[int(m.group(1))] = m.group(2).strip()
    return sections


def strip_md_inline(text: str) -> str:
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"(?<!\*)\*([^*\n]+)\*(?!\*)", r"\1", text)
    return text


def parse_md_table(block: str) -> list[list[str]]:
    rows: list[list[str]] = []
    for line in block.splitlines():
        line = line.strip()
        if not line.startswith("|"):
            continue
        if re.match(r"^\|[\s:|\-]+\|$", line):
            continue
        cells = [c.strip() for c in line.strip("|").split("|")]
        rows.append(cells)
    return rows


def table_to_list(block: str) -> str:
    rows = parse_md_table(block)
    if len(rows) < 2:
        return block
    data = rows[1:]  # drop header
    out: list[str] = []
    for row in data:
        first = row[0].strip()
        rest = [c for c in row[1:] if c]
        rest_text = " — ".join(rest)
        if first and ID_PATTERN.match(first):
            out.append(f"{first}. {rest_text}" if rest_text else f"{first}.")
        else:
            prefix = f"{first} — " if first else ""
            out.append(f"• {prefix}{rest_text}".rstrip(" —"))
    return "\n".join(out)


def replace_tables(text: str) -> str:
    lines = text.split("\n")
    result: list[str] = []
    i = 0
    while i < len(lines):
        is_table_start = (
            lines[i].lstrip().startswith("|")
            and i + 1 < len(lines)
            and re.match(r"^\s*\|[\s:|\-]+\|\s*$", lines[i + 1])
        )
        if is_table_start:
            block: list[str] = []
            while i < len(lines) and lines[i].lstrip().startswith("|"):
                block.append(lines[i])
                i += 1
            result.append(table_to_list("\n".join(block)))
        else:
            result.append(lines[i])
            i += 1
    return "\n".join(result)


def convert_bullets(text: str) -> str:
    out: list[str] = []
    for line in text.split("\n"):
        m = re.match(r"^(\s*)-\s+(.*)", line)
        if m:
            out.append(f"{m.group(1)}• {m.group(2)}")
        else:
            out.append(line)
    return "\n".join(out)


def normalise(text: str) -> str:
    text = replace_tables(text)
    text = convert_bullets(text)
    text = strip_md_inline(text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


EDGE_EXE = Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe")


def svg_to_png(svg_path: Path) -> Path:
    """Render SVG → PNG via headless Edge (writes a sibling .png file)."""
    import subprocess

    png_path = svg_path.with_suffix(".png")
    if png_path.exists() and png_path.stat().st_mtime >= svg_path.stat().st_mtime:
        return png_path

    text = svg_path.read_text(encoding="utf-8", errors="ignore")
    vb = re.search(r'viewBox=["\']([\d.\s\-]+)["\']', text)
    if vb:
        parts = vb.group(1).split()
        w, h = int(float(parts[2])), int(float(parts[3]))
    else:
        w, h = 1200, 900

    url = "file:///" + str(svg_path).replace("\\", "/")
    subprocess.run(
        [
            str(EDGE_EXE),
            "--headless=new",
            "--disable-gpu",
            f"--screenshot={png_path}",
            f"--window-size={w},{h}",
            "--hide-scrollbars",
            "--default-background-color=00000000",
            url,
        ],
        check=True,
        capture_output=True,
        timeout=60,
    )
    return png_path


def resolve_image(path: Path) -> Path:
    if path.suffix.lower() == ".svg":
        return svg_to_png(path)
    return path


def extract_images(content: str, base_dir: Path) -> tuple[str, list[Path]]:
    img_pattern = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")
    images: list[Path] = []
    for m in img_pattern.finditer(content):
        path = (base_dir / m.group(1)).resolve()
        if path.exists():
            try:
                images.append(resolve_image(path))
            except Exception as e:
                print(f"  ! couldn't process {path.name}: {e}", file=sys.stderr)
        else:
            print(f"  ! image not found: {path}", file=sys.stderr)
    cleaned = img_pattern.sub("", content)
    return cleaned, images


def build_slide(prs: Presentation, md_path: Path) -> None:
    sections = split_sections(md_path.read_text(encoding="utf-8"))

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    title_text = strip_md_inline(sections.get(1, "").strip())
    if title_text:
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.6))
        title_box.text_frame.text = title_text
        title_box.text_frame.word_wrap = True

    body_raw = sections.get(2, "")
    body_text, images = extract_images(body_raw, md_path.parent)
    body_text = normalise(body_text)

    if body_text:
        body_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(12.7), Inches(6.0))
        body_box.text_frame.text = body_text
        body_box.text_frame.word_wrap = True

    img_top = Inches(1.0)
    for img_path in images:
        try:
            slide.shapes.add_picture(str(img_path), Inches(0.3), img_top)
        except Exception as e:
            print(f"  ! failed to add {img_path.name}: {e}", file=sys.stderr)
        img_top += Inches(2)

    notes_parts: list[str] = []
    if sections.get(3):
        notes_parts.append(normalise(sections[3]))
    if sections.get(4):
        notes_parts.append("— Терминдер —")
        notes_parts.append(normalise(sections[4]))
    notes_text = "\n\n".join(notes_parts)
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text


def main(limit: int | None = None) -> None:
    md_files = sorted(SLIDES_DIR.glob("*.md"))
    if limit is not None:
        md_files = md_files[:limit]

    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU

    for md_path in md_files:
        print(f"Building: {md_path.name}")
        build_slide(prs, md_path)

    prs.save(str(OUTPUT_PATH))
    print(f"Saved: {OUTPUT_PATH} ({len(md_files)} slides)")


if __name__ == "__main__":
    n = int(sys.argv[1]) if len(sys.argv) > 1 else None
    main(limit=n)
